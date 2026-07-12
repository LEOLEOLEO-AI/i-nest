# -*- coding: utf-8 -*-
"""Import Genspark docs into Obsidian as TCC/iNEST first-principles MD files"""

from pathlib import Path
from docx import Document
from pptx import Presentation
import re, json

SRC = Path(r"D:\Output\Genspark")
DST = Path(r"D:\Obsidian\home\work\.openclaw\workspace")

def extract_docx(filepath):
    """Extract text from .docx, preserving structure"""
    doc = Document(filepath)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        style = para.style.name if para.style else ""
        if "Heading 1" in style or "heading 1" in style.lower():
            lines.append(f"# {text}")
        elif "Heading 2" in style or "heading 2" in style.lower():
            lines.append(f"## {text}")
        elif "Heading 3" in style or "heading 3" in style.lower():
            lines.append(f"### {text}")
        elif para.runs and para.runs[0].bold:
            lines.append(f"**{text}**")
        else:
            lines.append(text)
    return "\n\n".join(lines)

def extract_pptx(filepath):
    """Extract text from .pptx slides"""
    prs = Presentation(filepath)
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"## Slide {i+1}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
        lines.append("")
    return "\n\n".join(lines)

# === File mapping ===
files_map = {
    "TCC": {
        "tcc_paradigm_guide": {
            "src": "拓扑中心计算范式（TCC）导读：核心内涵、原语体系、业务映射与模式切换机理.docx",
            "title": "TCC第一性逻辑：核心内涵与架构映射",
            "dest": "30_TCC/31_Theory/tcc_first_principles.md"
        },
        "tcc_paper_background": {
            "src": "TCC拓扑中心计算范式：论文背景综述、相关工作对比与引用框架.docx",
            "title": "TCC论文背景综述与相关工作",
            "dest": "30_TCC/31_Theory/tcc_paper_background.md"
        },
        "tcc_pptx": {
            "src": "TCC拓扑中心计算范式_核心原理与架构映射_20260711091841.pptx",
            "title": "TCC核心原理与架构映射（PPT）",
            "dest": "30_TCC/31_Theory/tcc_core_principles_pptx.md"
        },
        "physical_limits": {
            "src": "The Physical Limits of Computing_ Why Topology Will Dominate the Post-Dennard Era.docx",
            "title": "物理极限：为何拓扑将主导后Dennard时代",
            "dest": "30_TCC/31_Theory/physical_limits_post_dennard.md"
        },
        "nonlinear_gain": {
            "src": "从节点堆砌到拓扑编程：系统级超非线性增益（1＋＞2）的理论基础、数学条件与产业证据.docx",
            "title": "从节点堆砌到拓扑编程：系统级超非线性增益",
            "dest": "30_TCC/31_Theory/nonlinear_gain_topology_programming.md"
        },
        "mesoscale": {
            "src": "以晶圆／晶矩为技术物理载体的介观尺度计算——“十五五”期间中国先进计算产业发展的思考与建议.docx",
            "title": "介观尺度计算：十五五先进计算发展思考",
            "dest": "30_TCC/34_Projects/mesoscale_computing_15th_five_year.md"
        },
    },
    "iNEST": {
        "inest_logic_chain": {
            "src": "iNEST的物理与生物角度推演过程——从复杂性第一性到网络时空协同复杂度智能涌现的逻辑链条复盘.docx",
            "title": "iNEST第一性逻辑：从复杂性到智能涌现的推演链条",
            "dest": "40_iNEST/41_Theory/inest_first_principles.md"
        },
        "autocatalytic": {
            "src": "自催化闭包与自催化网络：从化学网络理论到网络时空协同复杂度智能涌现的第一性原理.docx",
            "title": "自催化闭包：从化学网络到智能涌现的第一性原理",
            "dest": "40_iNEST/41_Theory/autocatalytic_closure_emergence.md"
        },
    }
}

stats = {"TCC": 0, "iNEST": 0, "errors": []}

for direction, items in files_map.items():
    for key, info in items.items():
        src_file = SRC / info["src"]
        if not src_file.exists():
            # Try alternate versions
            alt = list(SRC.glob(info["src"].replace(".docx", "*.docx").replace(".pptx", "*.pptx")))
            if alt:
                src_file = alt[0]
            else:
                stats["errors"].append(f"NOT FOUND: {info['src']}")
                continue
        
        print(f"[{direction}] {info['title'][:50]}...")
        
        try:
            if src_file.suffix.lower() == ".pptx":
                content = extract_pptx(src_file)
            else:
                content = extract_docx(src_file)
            
            # Add frontmatter and header
            header = f"""---
title: "{info['title']}"
direction: {direction}
source: "Genspark"
date: 2026-07-12
tags: [{direction.lower()}, first-principles, genspark-import]
---

# {info['title']}

> 来源: Genspark 创新引擎 | 方向: {direction} | 导入日期: 2026-07-12

---

"""
            full_content = header + content
            
            dest_file = DST / info["dest"]
            dest_file.parent.mkdir(parents=True, exist_ok=True)
            dest_file.write_text(full_content, encoding="utf-8")
            
            size_kb = dest_file.stat().st_size / 1024
            print(f"  -> {info['dest']} ({size_kb:.1f}KB)")
            stats[direction] += 1
            
        except Exception as e:
            stats["errors"].append(f"ERROR {info['src']}: {e}")
            print(f"  ERR: {e}")

print(f"\n=== Import Summary ===")
print(f"TCC files: {stats['TCC']}")
print(f"iNEST files: {stats['iNEST']}")
print(f"Errors: {len(stats['errors'])}")
for e in stats["errors"]:
    print(f"  {e}")
