# -*- coding: utf-8 -*-
"""
生成立项/评审用中文 PPT：
《耦合随机延迟网络的智能涌现阈值——面向物理类脑硬件的自由能视角》
规范：16:9；标题微软雅黑≥28磅；正文微软雅黑/Arial≥20磅；单页≤50字；每页一结论。
依赖：python-pptx  ( pip install python-pptx )
输出：iNEST_SDDE_涌现阈值.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------- 全局配色（学术深蓝主色） ----------------
NAVY   = RGBColor(0x1F, 0x4E, 0x79)   # 主色·深蓝
RED    = RGBColor(0xC0, 0x50, 0x4D)   # 强调·砖红
GRAY   = RGBColor(0x40, 0x40, 0x40)   # 正文·深灰
LIGHT  = RGBColor(0xF2, 0xF5, 0xFA)   # 背景·浅蓝
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "微软雅黑"
BODY_FONT  = "微软雅黑"
EN_FONT    = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)      # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------------- 通用函数 ----------------
def add_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_title(slide, text, size=30, color=NAVY, top=0.45):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.name = TITLE_FONT; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = color
    # 标题下划线条
    from pptx.enum.shapes import MSO_SHAPE
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(top+0.95),
                                Inches(3.2), Inches(0.06))
    ln.fill.solid(); ln.fill.fore_color.rgb = RED; ln.line.fill.background()
    return slide

def add_point(slide, text, top, size=22, color=GRAY, bold=False, left=1.0, width=11.3):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.name = BODY_FONT; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return tb

def add_conclusion_bar(slide, text, size=22):
    """页脚结论条：每页一句结论"""
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.7), Inches(6.35), Inches(12.0), Inches(0.75))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "结论：" + text
    r.font.name = BODY_FONT; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = WHITE

def add_formula(slide, text, top, size=26, color=NAVY):
    """公式以文本近似标准渲染（如需真渲染可后期粘贴图片）"""
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(11.3), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = "Cambria Math"; r.font.size = Pt(size); r.font.italic = True
    r.font.color.rgb = color
    return tb

# ================================================================
# 封面
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "耦合随机延迟网络的智能涌现阈值"
r.font.name = TITLE_FONT; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
r2 = p2.add_run(); r2.text = "——面向物理类脑硬件的自由能视角"
r2.font.name = TITLE_FONT; r2.font.size = Pt(26); r2.font.color.rgb = RGBColor(0xBD,0xD7,0xEE)
sub = s.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.5))
sf = sub.text_frame
rp = sf.paragraphs[0]; rr = rp.add_run()
rr.text = "网络时空协同复杂度涌现智能研究院（iNEST）"
rr.font.name = TITLE_FONT; rr.font.size = Pt(22); rr.font.color.rgb = WHITE

# ================================================================
# 第1页 研究背景
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "一、研究背景：真实系统既有记忆，又有涨落")
add_point(s, "确定性延迟方程（ＤＤＥ）刻画记忆", 2.0, 24, NAVY, True)
add_point(s, "随机常微分方程（ＳＯＤＥ）刻画涨落", 3.0, 24, NAVY, True)
add_point(s, "随机延迟微分方程（ＳＤＤＥ）统一二者", 4.0, 24, RED, True)
add_point(s, "Buckwar（2000）奠基单节点数值分析，被引 300+ 次", 5.2, 22, GRAY)
add_conclusion_bar(s, "记忆＋涨落须联立建模，ＳＤＤＥ 是物理类脑硬件的天然语言")

# ================================================================
# 第2页 科学问题
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "二、科学问题：复杂度何时越阈涌现智能？")
add_point(s, "iNEST 核心信念：时空复杂度超阈 → 智能涌现", 2.2, 24, NAVY, True)
add_point(s, "难点：从「哲学信念」到「可证定理」", 3.4, 24, RED, True)
add_point(s, "关键：找到支配相变的单一序参量", 4.6, 24, GRAY, True)
add_conclusion_bar(s, "把「复杂度越阈涌现」转化为可硬件验证的数学判据")

# ================================================================
# 第3页 模型
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "三、模型：单节点提升为异质耦合网络")
add_formula(s, "dX = F(X, Xτ) dt + G(X, Xτ) dW", 2.1, 28, NAVY)
add_point(s, "耦合矩阵 A ： 元拓扑与软件定义互连", 3.4, 22, GRAY, True)
add_point(s, "延迟矩阵 τ ： 异质时空互连延迟", 4.2, 22, GRAY, True)
add_point(s, "噪声矩阵 Σ ： 物理涨落空间相关", 5.0, 22, GRAY, True)
add_conclusion_bar(s, "三矩阵联合编码网络的时空协同复杂度")

# ================================================================
# 第4页 定理一：适定性与线性判据
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "四、定理一：适定性与涌现序参量")
add_point(s, "命题：强解存在唯一（步进法＋延迟 Grönwall）", 2.1, 22, NAVY, True)
add_formula(s, "Θ = λmax(B+Bᵀ) + 2μmax(A) + λg", 3.2, 26, NAVY)
add_point(s, "Θ < 0 ： 均方稳定（未涌现）", 4.4, 22, GRAY, True)
add_point(s, "Θ > Θc ： 随机分岔（涌现新序）", 5.2, 22, RED, True)
add_conclusion_bar(s, "单一序参量 Θ 支配「稳定—涌现」相变")

# ================================================================
# 第5页 定理二：尖锐非线性二分律
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "五、定理二：尖锐的非线性二分律")
add_point(s, "适用超线性节点（Khasminskii 单边条件）", 2.1, 22, NAVY, True)
add_formula(s, "Θ* = a₁ + a₂", 3.3, 28, RED)
add_point(s, "同一常数 Θ* 既定稳定，又定失稳", 4.5, 24, GRAY, True)
add_point(s, "充要闭环 → 达顶刊标准", 5.4, 22, NAVY, True)
add_conclusion_bar(s, "Θ* 是稳定与涌现的充要分界，判据尖锐")

# ================================================================
# 第6页 算例验证
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "六、算例：两节点 tanh 网络闭式验证")
add_formula(s, "Θ* = -2α + σ² + 2κ", 2.1, 28, NAVY)
add_point(s, "阈下 κ=0.3 ： Θ*=-1.24 → 指数衰减", 3.4, 22, GRAY, True)
add_point(s, "阈上 κ=1.3 ： Θ*=+0.76 → 增长后饱和", 4.4, 22, RED, True)
add_point(s, "蒙特卡洛经验率 0.86 对理论 0.88，吻合", 5.4, 22, NAVY)
add_conclusion_bar(s, "最小可验证实例双向印证定理二")

# ================================================================
# 第7页 自由能诠释与六级智能
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "七、诠释：Θ* 映射六级智能等级")
rows = [("Θ* ≪ 0", "感知 L1"), ("Θ* ≲ 0", "反应 L2"),
        ("0 < Θ* < Θc", "适应 L3"), ("Θ* ≈ Θc", "创造 L4"),
        ("Θ* > Θc", "通用/超级 L5/L6")]
top = 1.95
for cond, lvl in rows:
    add_point(s, cond, top, 22, NAVY, True, left=1.3, width=5.0)
    add_point(s, "→  " + lvl, top, 22, RED, True, left=7.0, width=5.0)
    top += 0.72
add_conclusion_bar(s, "延迟将分岔点展开为分岔带，抬升可达智能等级")

# ================================================================
# 第8页 价值与路线
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "八、价值与技术路线")
add_point(s, "理论：复杂度涌现从信念到定理", 2.1, 24, NAVY, True)
add_point(s, "硬件：晶圆级液态网络时序整数对齐", 3.2, 24, GRAY, True)
add_point(s, "范式：物理载体即可微延迟随机求解器", 4.3, 24, RED, True)
add_point(s, "驱动：最小自由能 + STDP 持续进化", 5.4, 22, NAVY)
add_conclusion_bar(s, "从算力堆砌智能迈向时空协同复杂度涌现智能")

# ================================================================
# 第9页 结语（封底）
# ================================================================
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "大道至简"
r.font.name = TITLE_FONT; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "让方程不再被求解，而被物理地演化出来"
r2.font.name = TITLE_FONT; r2.font.size = Pt(26); r2.font.color.rgb = RGBColor(0xBD,0xD7,0xEE)
note = s.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.6))
nr = note.text_frame.paragraphs[0].add_run()
nr.text = "本演示所有内容均基于公开资料研究"
nr.font.name = TITLE_FONT; nr.font.size = Pt(16); nr.font.color.rgb = RGBColor(0x9D,0xB8,0xD6)

# ---------------- 保存 ----------------
import os
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "iNEST_SDDE_涌现阈值.pptx")
prs.save(out_path)
print("已生成：", out_path, " 共", len(prs.slides._sldIdLst), "页")
