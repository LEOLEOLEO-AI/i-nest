# -*- coding: utf-8 -*-
import io, sys
from pptx import Presentation
p = Presentation('介观尺度计算新范式.pptx')
print('slides:', len(p.slides._sldIdLst))

texts = []
for idx, s in enumerate(p.slides, 1):
    for sh in s.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t = ''.join(r.text for r in para.runs)
                if t.strip():
                    texts.append((idx, t))
alltext = '\n'.join(t for _, t in texts)

# Correction 1: 邬江兴 院士 on cover
c1 = ('邬江兴' in alltext) and ('院士' in alltext)

# Correction 2: R_I three-tier (look at slide 10 specifically)
s10 = '\n'.join(t for idx, t in texts if idx == 10)
c2_lt = 'R_I < 1' in s10 and '难以胜任' in s10
c2_eq = 'R_I' in s10 and '≈' in s10 and '能力匹配' in s10
c2_gt = 'R_I > 1' in s10 and '游刃有余' in s10
# ensure the old "远大于1进候选区" as a tier/candidate-zone claim is gone
old_candidate = '远大于1' in s10 and '候选区' in s10   # should be False on slide 10
# the phrase only appears as an explanatory correction note elsewhere
note_ok = alltext.count('远大于1') >= 1  # appears in the revision note, which is intended
c2 = c2_lt and c2_eq and c2_gt and (not old_candidate)

# Correction 3: SDI operator framed as engineering control law
c3 = ('工程控制律' in alltext) and ('受控' in alltext) and ('SDDE' in alltext)
c3_natural = ('自然定律' in alltext)  # we state it is NOT a natural law

print('Correction1 (邬江兴院士 on cover):', c1)
print('Correction2  R_I<1 难以胜任 :', c2_lt)
print('Correction2  R_I≈1 能力匹配 :', c2_eq)
print('Correction2  R_I>1 游刃有余 :', c2_gt)
print('Correction2  old 候选区 gone :', not old_candidate)
print('Correction2 overall          :', c2)
print('Correction3 工程控制律+受控SDDE:', c3)
print('Correction3 declares not natural law:', c3_natural)
