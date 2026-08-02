# -*- coding: utf-8 -*-
"""
Build the 20-page academic deck: 介观尺度计算新范式——从SDSoW到iNEST智能涌现之路
Style: 学术风 (Academic) — deep blue / white / light blue, formula-first.
Three key corrections applied:
  1. Cover reports 邬江兴院士.
  2. R_I criterion: <1 难以胜任 / ≈1 能力匹配 / >1 游刃有余 (no longer "远大于1进候选区").
  3. A(t^+)=Π_SDI[...] is the SDI topology-control operator — an engineering control law
     abstracted from the controlled SDDE hybrid dynamics system, not a natural-law formula.
"""
import os, math
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- palette ----------------
NAVY  = RGBColor(0x0E, 0x3F, 0x8C)
BLUE  = RGBColor(0x1E, 0x4F, 0xA8)
LBLUE = RGBColor(0xE8, 0xEF, 0xF8)
LBLUE2= RGBColor(0xF0, 0xF5, 0xFC)
MIDB  = RGBColor(0x3D, 0x7B, 0xD9)
RED   = RGBColor(0xD9, 0x53, 0x4F)
YEL   = RGBColor(0xFF, 0xC1, 0x07)
YELBG = RGBColor(0xFF, 0xF6, 0xE6)
BLUEBG= RGBColor(0xE8, 0xEF, 0xF8)
TEXT  = RGBColor(0x1A, 0x22, 0x30)
SUB   = RGBColor(0x4A, 0x55, 0x68)
AUX   = RGBColor(0x8B, 0x97, 0xA8)
LINE  = RGBColor(0xD6, 0xDC, 0xE5)
LINE2 = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE  = RGBColor(0xEE, 0xF1, 0xF5)

FONT = "Microsoft YaHei"
SERIF = "Times New Roman"

SW, SH = 13.333, 7.5
TITLE_H = 0.85
FOOTER_Y = 6.95
CTOP = 1.0
CBOT = 6.85

OUTDIR = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(OUTDIR, "resources", "images")
os.makedirs(IMG, exist_ok=True)

# ---------------- formula rendering ----------------
_formula_cache = {}
def render_formula(latex, name, fontsize=32, color="#0E3F8C", dpi=300):
    path = os.path.join(IMG, name + ".png")
    if name in _formula_cache:
        return _formula_cache[name]
    fig = plt.figure(figsize=(0.1, 0.1), dpi=dpi)
    fig.patch.set_alpha(0.0)
    fig.text(0.5, 0.5, "$" + latex + "$", fontsize=fontsize, color=color,
             ha="center", va="center")
    fig.savefig(path, dpi=dpi, transparent=True, bbox_inches="tight", pad_inches=0.03)
    plt.close(fig)
    _formula_cache[name] = path
    return path

# ---------------- low-level helpers ----------------
def set_run_font(run, name=FONT, size=18, color=TEXT, bold=False, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=True, radius=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75 * line_w)
    shp.shadow.inherit = False
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp

def add_text(slide, l, t, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, fill=None, line=None, rounded=True, space_after=4):
    """paras: list of paragraphs; each paragraph is list of (txt,size,color,bold,italic)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    if fill is not None or line is not None:
        # wrap textbox in a shape background
        shp = add_rect(slide, l, t, w, h, fill=fill, line=line, rounded=rounded)
        shp.shadow.inherit = False
        # textbox on top (transparent)
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        for (txt, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt
            set_run_font(r, FONT, size, color, bold, italic)
    return tb

def add_pic(slide, path, l, t, width=None, height=None):
    kw = {}
    if width is not None: kw["width"] = Inches(width)
    if height is not None: kw["height"] = Inches(height)
    return slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def header(slide, title, tag):
    add_rect(slide, 0, 0, SW, TITLE_H, fill=NAVY, rounded=False)
    add_text(slide, 0.55, 0.0, SW - 2.6, TITLE_H,
             [[(title, 24, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - 2.6, 0.0, 2.1, TITLE_H,
             [[(tag, 12, RGBColor(0xC9, 0xD8, 0xF0), False, False)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def footer(slide, num):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(FOOTER_Y), Inches(SW), Pt(1.0))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE2; ln.line.fill.background(); ln.shadow.inherit = False
    add_text(slide, 0.55, FOOTER_Y + 0.02, 9.0, 0.5,
             [[("介观尺度计算新范式 · 从SDSoW到iNEST", 10.5, AUX, False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - 2.0, FOOTER_Y + 0.02, 1.45, 0.5,
             [[(f"{num:02d} / 20", 10.5, BLUE, True, False)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_table(slide, l, t, w, h, data, col_widths=None, header_fill=NAVY,
              font_size=14, hdr_size=15):
    rows = len(data); cols = len(data[0])
    gf = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = gf.table
    if col_widths:
        tot = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(w * cw / tot)
    tbl.first_row = False; tbl.horz_banding = False
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            val = data[r][c]
            if isinstance(val, tuple):
                txt, color, bold = val
            else:
                txt = val; bold = (r == 0)
                color = WHITE if r == 0 else TEXT
            run = para.add_run(); run.text = txt
            set_run_font(run, FONT, hdr_size if r == 0 else font_size, color, bold, False)
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = (LBLUE2 if r % 2 == 1 else WHITE)
    return tbl

# ============================================================
# SLIDE BUILDERS
# ============================================================
def s01_cover(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SW, SH, fill=NAVY, rounded=False)
    # accent strip
    add_rect(s, 0, 0, 0.35, SH, fill=MIDB, rounded=False)
    add_text(s, 0.8, 0.55, 11.7, 0.5,
             [[("国家数字交换系统工程技术研究中心 · 复旦大学大数据研究院", 14, RGBColor(0xC9,0xD8,0xF0), False, False)]])
    add_text(s, SW - 3.5, 0.5, 2.9, 0.5,
             [[("后冯·诺依曼时代 · 换道引领战略机遇", 12, RGBColor(0x9F,0xB6,0xE0), False, False)]],
             align=PP_ALIGN.RIGHT)
    add_text(s, 0.8, 1.7, 11.7, 0.4,
             [[("MESOSCOPIC COMPUTING NEW PARADIGM", 14, RGBColor(0x9F,0xB6,0xE0), False, True)]])
    add_text(s, 0.8, 2.15, 11.7, 1.0,
             [[("介观尺度计算新范式", 46, WHITE, True, False)]])
    add_text(s, 0.8, 3.2, 11.7, 0.7,
             [[("——从 SDSoW 到 iNEST 智能涌现之路", 28, RGBColor(0xEA,0xF1,0xFB), True, False)]])
    # evolution chain
    nodes = [("晶圆 / 晶矩平台", "介观物理资源"), ("液态拓扑", "复杂度匹配"),
             ("智涌脑", "智能涌现"), ("具身智能体", "自主行为")]
    x = 0.8; y = 4.5; bw = 2.55; bh = 1.0; gap = 0.45
    for i, (t1, t2) in enumerate(nodes):
        add_rect(s, x, y, bw, bh, fill=RGBColor(0x16,0x4A,0x9C), line=MIDB, line_w=1.2)
        add_text(s, x, y + 0.12, bw, 0.5, [[(t1, 17, WHITE, True, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + 0.55, bw, 0.35, [[(t2, 12, RGBColor(0xBB,0xD0,0xF0), False, False)]],
                 align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, x + bw, y, gap, bh, [[("→", 22, YEL, True, False)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += bw + gap
    # bottom: reporter + quote
    add_text(s, 0.8, 6.1, 6.0, 0.5,
             [[("报告人：", 16, RGBColor(0xDC, 0xE7, 0xF8), False, False),
               ("邬江兴 院士", 18, WHITE, True, False)]])
    add_text(s, 0.8, 6.55, 6.5, 0.5,
             [[("国家数字交换系统工程技术研究中心 · 复旦大学大数据研究院", 12, RGBColor(0xC9,0xD8,0xF0), False, False)]])
    add_rect(s, 7.6, 6.05, 5.1, 1.05, fill=RGBColor(0x16,0x4A,0x9C), line=YEL, line_w=1.5, rounded=False)
    add_text(s, 7.85, 6.2, 4.6, 0.8,
             [[("“让物理网络自己长出智能”", 19, YEL, True, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def s02_frontier(prs):
    s = blank(prs); header(s, "从国家前沿科学问题说起", "引子 · 02")
    add_rect(s, 0.55, CTOP + 0.1, 4.0, CBOT - CTOP - 0.2, fill=NAVY, rounded=False)
    add_text(s, 0.85, CTOP + 0.4, 3.4, 0.4, [[("FRONTIER QUESTION", 13, RGBColor(0x9F,0xB6,0xE0), False, False)]])
    add_text(s, 0.85, CTOP + 1.0, 3.4, 1.6,
             [[("类生物", 30, WHITE, True, False)], [("计算", 30, YEL, True, False)], [("成为命题", 30, WHITE, True, False)]])
    add_rect(s, 0.85, CBOT - 1.7, 3.4, 1.4, fill=RGBColor(0x16,0x4A,0x9C), line=None, rounded=True)
    add_text(s, 1.05, CBOT - 1.55, 3.0, 1.1,
             [[("智能涌现的", 13, RGBColor(0xDC,0xE7,0xF8), False, False),
               ("第一性原理", 14, WHITE, True, False),
               ("，是介观尺度的核心科学之问。", 13, RGBColor(0xDC,0xE7,0xF8), False, False)]])
    lx = 4.95
    add_text(s, lx, CTOP + 0.2, 7.7, 0.6,
             [[("类生物计算，正在成为前沿科学命题", 22, BLUE, True, False)]])
    add_rect(s, lx, CTOP + 0.85, 7.7, 0.03, fill=LINE, rounded=False)
    add_text(s, lx, CTOP + 1.1, 7.7, 1.8,
             [[("当前，人工智能的瓶颈已不再是单纯的数据规模或训练技巧，而是对“智能从何而来”的根本追问。类生物计算的关键，不只是仿脑结构，而是要回答：",
                17, TEXT, False, False),
               ("智能涌现的第一性原理是什么？它能否在介观尺度上被定义、被度量、被构建？",
                17, NAVY, True, False)]], space_after=8)
    add_text(s, lx, CTOP + 3.3, 7.7, 1.2,
             [[("这把问题从“造更快的机器”转向“造会生长的机器”。当物理网络具备可塑结构与非线性动力学，智能不再只是被编程的结果，而可能成为系统自身演化的涌现属性——这正是本报告的起点。",
                17, SUB, False, False)]])
    add_rect(s, lx, CBOT - 1.0, 7.7, 0.8, fill=LBLUE2, line=BLUE, line_w=1.0)
    add_text(s, lx + 0.2, CBOT - 0.92, 7.4, 0.65,
             [[("页脚佐证：中国科协公开发布信息；IT之家公开报道，2026。", 12, AUX, False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 2)

def s03_walls(prs):
    s = blank(prs); header(s, "不是算力不够，而是范式太旧", "引子 · 03")
    add_rect(s, 0.55, CTOP + 0.1, 4.3, CBOT - CTOP - 0.2, fill=NAVY, rounded=False)
    add_text(s, 0.9, CTOP + 0.5, 3.6, 1.4,
             [[("3", 90, YEL, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.9, CTOP + 2.2, 3.6, 0.6, [[("堵墙同时逼近", 22, WHITE, True, False)]])
    add_text(s, 0.9, CTOP + 2.9, 3.6, 1.2,
             [[("功耗墙 · 存储墙 · 互连墙，正把“堆算力”路线的红利吃干。", 15, RGBColor(0xC9,0xD8,0xF0), False, False)]])
    add_rect(s, 0.9, CBOT - 1.7, 3.6, 1.4, fill=RGBColor(0x16,0x4A,0x9C), rounded=True)
    add_text(s, 1.1, CBOT - 1.55, 3.2, 1.1,
             [[("主矛盾已从", 14, RGBColor(0xEA,0xF1,0xFB), False, False),
               ("单点算力", 15, YEL, True, False),
               ("，转向", 14, RGBColor(0xEA,0xF1,0xFB), False, False),
               ("系统架构", 15, YEL, True, False), ("。", 14, RGBColor(0xEA,0xF1,0xFB), False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    walls = [("功耗墙", "Power Wall", "AI 算力堆叠带来能耗陡增，单位算力的能效正逼近半导体物理下限，单纯扩大规模难以为继。"),
             ("存储墙", "Memory Wall", "数据搬运成为时延与能耗的主要来源，大量算力“饿死”在访存等待中，算力利用率被严重拖累。"),
             ("互连墙", "Interconnect Wall", "系统规模扩张受限于互连带宽、拓扑扩展与同步开销，横向堆核遭遇通信瓶颈。")]
    lx = 5.25; wy = CTOP + 0.15; wh = (CBOT - CTOP - 0.6) / 3
    for i, (k, e, d) in enumerate(walls):
        yy = wy + i * (wh + 0.18)
        add_rect(s, lx, yy, 7.45, wh, fill=LBLUE2, line=LINE, line_w=1.0)
        add_rect(s, lx, yy, 0.12, wh, fill=BLUE, rounded=False)
        add_text(s, lx + 0.3, yy + 0.12, 7.0, 0.5,
                 [[(k, 19, NAVY, True, False), ("   " + e, 13, AUX, False, False)]])
        add_text(s, lx + 0.3, yy + 0.62, 7.0, wh - 0.7,
                 [[(d, 15, SUB, False, False)]])
    footer(s, 3)

def s04_moore(prs):
    s = blank(prs); header(s, "系统级摩尔定律正在兴起", "引子 · 04")
    # left diagram
    add_rect(s, 0.55, CTOP + 0.1, 6.6, CBOT - CTOP - 0.2, fill=LBLUE2, line=LINE, line_w=1.0)
    add_text(s, 0.85, CTOP + 0.35, 6.0, 0.4,
             [[("竞争焦点迁移：从微缩到集成", 16, NAVY, True, False)]], align=PP_ALIGN.CENTER)
    add_rect(s, 1.6, CTOP + 1.1, 4.5, 1.0, fill=LBLUE, line=MIDB, line_w=1.0, rounded=True)
    add_text(s, 1.6, CTOP + 1.1, 4.5, 1.0,
             [[("旧范式：先进节点微缩", 16, SUB, False, False)],
              [("逼近物理极限，红利衰减", 12, AUX, False, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.6, CTOP + 2.2, 4.5, 0.5, [[("↓", 24, BLUE, True, False)]], align=PP_ALIGN.CENTER)
    add_rect(s, 1.6, CTOP + 2.75, 4.5, 0.95, fill=NAVY, rounded=True)
    add_text(s, 1.6, CTOP + 2.75, 4.5, 0.95,
             [[("新范式：系统级摩尔定律", 18, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    subs = ["SoW\n晶圆级", "先进封装\n2.5D / 3D", "3DHI\n异质集成"]
    bw = 1.95; gap = 0.2; x0 = 0.95
    for i, t in enumerate(subs):
        add_rect(s, x0 + i * (bw + gap), CTOP + 4.0, bw, 1.1, fill=WHITE, line=MIDB, line_w=1.2, rounded=True)
        add_text(s, x0 + i * (bw + gap), CTOP + 4.0, bw, 1.1,
                 [[(t.replace("\n", "\n"), 16, BLUE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # right text
    lx = 7.45
    add_text(s, lx, CTOP + 0.3, 5.3, 0.6, [[("SoW、先进封装、3DHI 成为主航道", 20, BLUE, True, False)]])
    add_rect(s, lx, CTOP + 0.95, 5.3, 0.03, fill=LINE, rounded=False)
    add_text(s, lx, CTOP + 1.2, 5.3, 2.2,
             [[("全球竞争的主战场，正从“把晶体管做得更小”，转向“把系统集得更巧”：晶圆级集成（SoW）、Chiplet 先进封装、三维异质集成（3DHI）共同构成",
                16, TEXT, False, False), ("系统级摩尔定律", 16, NAVY, True, False), ("。", 16, TEXT, False, False)]], space_after=8)
    add_text(s, lx, CTOP + 3.6, 5.3, 1.4,
             [[("这意味着价值重心上移到体系结构与互连创新——也为介观物理网络提供了难得的换道窗口。", 16, SUB, False, False)]])
    add_rect(s, lx, CBOT - 1.0, 5.3, 0.8, fill=LBLUE2, line=BLUE, line_w=1.0)
    add_text(s, lx + 0.2, CBOT - 0.92, 5.0, 0.65,
             [[("页脚佐证：IRDS；Heterogeneous Integration Roadmap；DARPA NGMM；NSTC；TSMC SoW；Cerebras WSE。", 11, AUX, False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 4)

def s05_sdsow(prs):
    s = blank(prs); header(s, "SDSoW 让介观网络资源可定义", "资源与刻度 · 05")
    add_rect(s, 0.55, CTOP + 0.1, 5.4, CBOT - CTOP - 0.2, fill=LBLUE2, line=LINE, line_w=1.0)
    add_text(s, 0.85, CTOP + 0.35, 4.8, 0.4,
             [[("资源尺度轴（自下而上扩展）", 15, NAVY, True, False)]], align=PP_ALIGN.CENTER)
    levels = [("芯片 / Chiplet", 0.62, LBLUE), ("晶圆 / 晶矩", 0.76, RGBColor(0xD6,0xE4,0xF7)),
              ("面板级", 0.88, RGBColor(0xBB,0xD2,0xF2)), ("机架 / 集群", 1.0, NAVY)]
    yy = CTOP + 1.0
    for t, frac, col in levels:
        bw = 4.6 * frac + 0.4
        add_rect(s, 0.85 + (5.0 - bw) / 2, yy, bw, 0.78, fill=col, rounded=True)
        add_text(s, 0.85 + (5.0 - bw) / 2, yy, bw, 0.78,
                 [[(t, 16, WHITE if col == NAVY else NAVY, True, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if t != "机架 / 集群":
            add_text(s, 0.85 + (5.0 - bw) / 2 + bw, yy, 0.5, 0.78, [[("↑", 18, MIDB, True, False)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.02
    lx = 6.35
    add_text(s, lx, CTOP + 0.3, 6.4, 0.6, [[("高密度 · 高维度 · 大规模 · 动态可塑", 20, BLUE, True, False)]])
    add_rect(s, lx, CTOP + 0.95, 6.4, 0.03, fill=LINE, rounded=False)
    add_text(s, lx, CTOP + 1.2, 6.4, 1.6,
             [[("SDSoW（Software-Defined System on Wafer）在", 16, TEXT, False, False),
               ("晶圆 / 晶矩 / 面板级", 16, NAVY, True, False),
               ("尺度上，把物理资源组织成可重构网络，兼具芯片的密度与系统的规模。", 16, TEXT, False, False)]], space_after=8)
    add_text(s, lx, CTOP + 3.2, 6.4, 1.6,
             [[("它首次让“介观物理网络资源”成为可被定义、调度与编程的对象——这是后续所有复杂度调控的物质前提。", 16, SUB, False, False)]])
    add_rect(s, lx, CBOT - 1.0, 6.4, 0.8, fill=LBLUE2, line=BLUE, line_w=1.0)
    add_text(s, lx + 0.2, CBOT - 0.92, 6.1, 0.65,
             [[("页脚佐证：IRDS；Heterogeneous Integration Roadmap；Cerebras WSE 公开资料。", 11, AUX, False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 5)

def s06_nonlinear(prs):
    s = blank(prs); header(s, "有资源，还要有非线性增益", "资源与刻度 · 06")
    add_text(s, 0, CTOP + 0.3, SW, 0.5,
             [[("关键是把结构资源转化为系统级非线性", 24, BLUE, True, False)]], align=PP_ALIGN.CENTER)
    add_text(s, 0, CTOP + 1.4, SW, 1.6,
             [[("1 + 1 > N", 110, NAVY, True, False)]], align=PP_ALIGN.CENTER)
    add_text(s, 1.5, CTOP + 3.4, 10.3, 1.6,
             [[("高密度资源只是基础。真正的问题是：如何让这些资源通过", 19, TEXT, False, False),
               ("结构组织", 19, NAVY, True, False),
               ("产生 1+1>2 乃至 1+1>N 的能力跃迁？", 19, TEXT, False, False)],
              [("这要求从“堆资源”转向“组结构”——用拓扑与动力学唤醒资源的非线性协同。", 18, SUB, False, False)]],
             align=PP_ALIGN.CENTER, space_after=10)
    add_text(s, 0, CBOT - 0.5, SW, 0.4,
             [[("页脚佐证：Anderson, Science, 1972（More is different）。", 12, AUX, False, False)]], align=PP_ALIGN.CENTER)
    footer(s, 6)

def s07_threshold(prs):
    s = blank(prs); header(s, "复杂度要从不可说走向可度量", "资源与刻度 · 07")
    add_rect(s, 0.55, CTOP + 0.1, 4.0, CBOT - CTOP - 0.2, fill=NAVY, rounded=False)
    add_text(s, 0.85, CTOP + 0.4, 3.4, 0.4, [[("THRESHOLD → CST", 13, RGBColor(0x9F,0xB6,0xE0), False, False)]])
    add_text(s, 0.85, CTOP + 1.0, 3.4, 1.6,
             [[("先立刻度", 28, WHITE, True, False)], [("再谈涌现", 28, WHITE, True, False)]])
    add_rect(s, 0.85, CBOT - 2.0, 3.4, 1.7, fill=RGBColor(0x16,0x4A,0x9C), line=YEL, line_w=1.2, rounded=True)
    add_text(s, 1.05, CBOT - 1.85, 3.0, 1.4,
             [[("“复杂自动机的组织，存在使功能跃迁的阈值。”", 14, RGBColor(0xEA,0xF1,0xFB), False, True)],
              [("—— 借鉴冯·诺依曼阈值思想的问题起点", 12, RGBColor(0x9F,0xB6,0xE0), False, False)]])
    lx = 4.95
    add_text(s, lx, CTOP + 0.2, 7.7, 0.6, [[("先立复杂度刻度，再谈智能涌现", 22, BLUE, True, False)]])
    add_rect(s, lx, CTOP + 0.85, 7.7, 0.03, fill=LINE, rounded=False)
    add_text(s, lx, CTOP + 1.1, 7.7, 1.9,
             [[("借鉴冯·诺依曼关于复杂自动机组织的", 17, TEXT, False, False),
               ("阈值思想", 17, NAVY, True, False),
               ("，我们并不宣称他提出了 CST 公式，而是以其思想为问题起点：首先定义一套复杂度刻度 CST，使系统的复杂度变得",
                17, TEXT, False, False), ("可观测、可度量、可预测、可调控", 17, BLUE, True, False), ("。", 17, TEXT, False, False)]], space_after=8)
    add_text(s, lx, CTOP + 3.4, 7.7, 1.3,
             [[("没有刻度，智能涌现就停留在比喻层面；有了刻度，匹配、同步与涌现才能成为可被工程化追求的目标。", 17, SUB, False, False)]])
    add_rect(s, lx, CBOT - 1.15, 7.7, 0.95, fill=YELBG, line=YEL, line_w=1.2)
    add_text(s, lx + 0.2, CBOT - 1.07, 7.4, 0.8,
             [[("口播边界：不说“冯·诺依曼提出了 CST 公式”，仅说其阈值思想启发了问题起点。", 14, RGBColor(0x8A,0x6D,0x2B), False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, lx, CBOT - 0.05, 7.7, 0.4,
             [[("页脚佐证：von Neumann, 1956；von Neumann, 1966；Ashby, 1956。", 11, AUX, False, False)]])
    footer(s, 7)

def s08_cst(prs):
    s = blank(prs); header(s, "可用时空协同复杂度 · CST V4.0", "资源与刻度 · 08")
    add_text(s, 0, CTOP + 0.15, SW, 0.4,
             [[("COMPLEXITY OF SPACE–TIME SYNERGY (CAPABLE)", 13, AUX, False, False)]], align=PP_ALIGN.CENTER)
    fp = render_formula(
        r"C_{\mathrm{ST}}^{\mathrm{cap}}(t) = S_c(t)\cdot T_c(t)\cdot \exp\left[ \alpha_{\mathrm{eff}}(t)\cdot \Gamma_{\mathrm{st}}^{u}(t) \right]",
        "cst", fontsize=34, color="#0E3F8C")
    add_rect(s, 1.4, CTOP + 0.7, 10.5, 1.7, fill=LBLUE2, line=MIDB, line_w=1.2)
    add_pic(s, fp, 1.9, CTOP + 0.85, width=9.5)
    params = [("S_c", "空间拓扑组织能力"), ("T_c", "时间演化与记忆能力"),
              ("α_eff", "非平衡态有效状态容量"), ("Γ_st^u", "结构·功能·环境匹配度")]
    bw = 2.85; gap = 0.25; x0 = 1.4
    for i, (k, d) in enumerate(params):
        x = x0 + i * (bw + gap)
        add_rect(s, x, CTOP + 2.7, bw, 1.0, fill=WHITE, line=LINE, line_w=1.0)
        fp2 = render_formula(k, "p_" + k.replace("/", "_").replace("^", ""), fontsize=22, color="#1E4FA8")
        add_pic(s, fp2, x + bw/2 - 0.6, CTOP + 2.78, width=1.2)
        add_text(s, x, CTOP + 3.35, bw, 0.3, [[(d, 12, SUB, False, False)]], align=PP_ALIGN.CENTER)
    add_rect(s, 1.4, CBOT - 1.05, 10.5, 0.9, fill=YELBG, line=YEL, line_w=1.2)
    add_text(s, 1.6, CBOT - 0.97, 10.1, 0.75,
             [[("金句：S_c·T_c 是复杂度底座，指数项是协同放大器。", 18, RGBColor(0x8A,0x6D,0x2B), True, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0, CBOT + 0.02, SW, 0.35,
             [[("页脚佐证：Anderson, 1972；Watts & Strogatz, Nature, 1998；Bullmore & Sporns, NRN, 2009。", 11, AUX, False, False)]],
             align=PP_ALIGN.CENTER)
    footer(s, 8)

def s09_exponent(prs):
    s = blank(prs); header(s, "指数放大从材料与非线性动力学中来", "资源与刻度 · 09")
    lx = 0.55
    add_text(s, lx, CTOP + 0.2, 7.0, 0.6, [[("材料状态容量 × 时空协同匹配", 20, BLUE, True, False)]])
    add_rect(s, lx, CTOP + 0.85, 7.0, 0.03, fill=LINE, rounded=False)
    add_text(s, lx, CTOP + 1.1, 7.0, 1.8,
             [[("Mott 相变、忆阻、相变、铁电、离子迁移等", 16, TEXT, False, False),
               ("非平衡过程", 16, NAVY, True, False),
               ("为系统提供内部状态变量。当材料状态容量与网络拓扑、时间动力学和环境需求匹配时，系统响应进入协同放大区。", 16, TEXT, False, False)]], space_after=8)
    fp1 = render_formula(r"\alpha_{\mathrm{eff}} \approx \ln M_{\mathrm{eff}}", "alpha_eff", fontsize=24, color="#0E3F8C")
    fp2 = render_formula(r"I(z_i; y_i)", "mutual_info", fontsize=24, color="#0E3F8C")
    add_rect(s, lx, CTOP + 3.2, 7.0, 1.1, fill=LBLUE2, line=MIDB, line_w=1.0)
    add_text(s, lx + 0.2, CTOP + 3.3, 3.3, 0.95, [[("有效状态容量：", 15, TEXT, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    add_pic(s, fp1, lx + 3.0, CTOP + 3.35, width=2.0)
    add_pic(s, fp2, lx + 5.1, CTOP + 3.45, width=1.7)
    add_text(s, lx, CBOT - 0.5, 7.0, 0.4,
             [[("页脚佐证：Chua, 1971；Strukov et al., Nature, 2008；Pickett et al., Nature Materials, 2013；Marković et al., Nat Rev Phys, 2020。", 10.5, AUX, False, False)]])
    # right material family
    add_rect(s, 7.85, CTOP + 0.1, 4.9, CBOT - CTOP - 0.2, fill=LBLUE2, line=LINE, line_w=1.0)
    add_text(s, 8.05, CTOP + 0.35, 4.5, 0.4,
             [[("内部状态变量来源（非线性器件）", 14, NAVY, True, False)]], align=PP_ALIGN.CENTER)
    mats = ["Mott 相变", "忆阻", "相变", "铁电", "离子迁移", "非线性动力学"]
    bw = 2.25; bh = 0.85; gap = 0.25; x0 = 8.15; y0 = CTOP + 1.1
    for i, m in enumerate(mats):
        r = i // 2; c = i % 2
        add_rect(s, x0 + c * (bw + gap), y0 + r * (bh + 0.3), bw, bh, fill=WHITE, line=MIDB, line_w=1.2, rounded=True)
        add_text(s, x0 + c * (bw + gap), y0 + r * (bh + 0.3), bw, bh,
                 [[(m, 16, BLUE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 8.15, CBOT - 1.15, 4.3, 0.85, fill=NAVY, rounded=True)
    add_text(s, 8.15, CBOT - 1.15, 4.3, 0.85,
             [[("容量 × 拓扑 × 动力学 × 环境 → 协同放大区", 14, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 9)

def s10_ri(prs):
    s = blank(prs); header(s, "智能是系统复杂度对环境复杂度的相对裕度", "度量与匹配 · 10")
    add_text(s, 0, CTOP + 0.12, SW, 0.4,
             [[("RELATIVE INTELLIGENCE INDEX", 13, AUX, False, False)]], align=PP_ALIGN.CENTER)
    fp = render_formula(r"R_I(t) = \frac{C_{\mathrm{ST}}^{\mathrm{sys}}(t)}{C_{\mathrm{ST}}^{\mathrm{env}}(t)}", "ri", fontsize=36, color="#0E3F8C")
    add_rect(s, 2.4, CTOP + 0.65, 8.5, 1.6, fill=LBLUE2, line=MIDB, line_w=1.2)
    add_pic(s, fp, 3.4, CTOP + 0.8, width=6.5)
    add_text(s, 1.0, CTOP + 2.4, 11.3, 0.5,
             [[("C_ST^sys", 15, NAVY, True, False), ("：系统可用时空协同复杂度　·　", 15, SUB, False, False),
               ("C_ST^env", 15, NAVY, True, False),
               ("：广义环境复杂度（环境状态·任务目标·扰动强度·实时约束·行动边界）", 15, SUB, False, False)]], align=PP_ALIGN.CENTER)
    tiers = [("R_I < 1", "难以胜任", PALE, RED),
             ("R_I ≈ 1", "能力匹配", LBLUE, BLUE),
             ("R_I > 1", "游刃有余", RGBColor(0xD6,0xE4,0xF7), NAVY)]
    bw = 3.7; gap = 0.35; x0 = 1.4
    for i, (k, d, bg, fg) in enumerate(tiers):
        x = x0 + i * (bw + gap)
        add_rect(s, x, CTOP + 3.2, bw, 1.15, fill=bg, line=fg, line_w=1.2)
        add_text(s, x, CTOP + 3.3, bw, 0.55, [[(k, 22, fg, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, CTOP + 3.85, bw, 0.45, [[(d, 14, SUB, False, False)]], align=PP_ALIGN.CENTER)
    add_rect(s, 1.4, CBOT - 1.05, 10.5, 0.9, fill=YELBG, line=YEL, line_w=1.2)
    add_text(s, 1.6, CBOT - 0.97, 10.1, 0.75,
             [[("判据修订：不再以单一阈值“进候选区”论定，统一为 小于1 / 约等于1 / 大于1 三档；不同区间对应不同智能等级（见 P11）。", 14, RGBColor(0x8A,0x6D,0x2B), False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0, CBOT + 0.02, SW, 0.35,
             [[("页脚佐证：Ashby, 1956；Friston, Nature Reviews Neuroscience, 2010。", 11, AUX, False, False)]], align=PP_ALIGN.CENTER)
    footer(s, 10)

def s11_levels(prs):
    s = blank(prs); header(s, "不同智能裕度，对应不同智能等级", "度量与匹配 · 11")
    add_rect(s, 0.55, CTOP + 0.1, 4.0, CBOT - CTOP - 0.2, fill=NAVY, rounded=False)
    add_text(s, 0.85, CTOP + 0.4, 3.4, 0.4, [[("SIX-LEVEL FRAMEWORK", 13, RGBColor(0x9F,0xB6,0xE0), False, False)]])
    add_text(s, 0.85, CTOP + 1.0, 3.4, 1.6, [[("裕度决定", 28, WHITE, True, False)], [("智能等级", 28, WHITE, True, False)]])
    fp = render_formula(r"L = \max_k\left\{ R_I^{(k)} \geq 1 \right\}", "level", fontsize=20, color="#FFFFFF")
    add_rect(s, 0.85, CBOT - 1.9, 3.4, 1.6, fill=RGBColor(0x16,0x4A,0x9C), rounded=True)
    add_text(s, 0.95, CBOT - 1.8, 3.2, 0.5, [[("等级由最大可胜任复杂度决定：", 12, RGBColor(0xDC,0xE7,0xF8), False, False)]])
    add_pic(s, fp, 1.0, CBOT - 1.3, width=3.1)
    levels = [("L1", "感知", "提取稳定环境信息", LBLUE),
              ("L2", "反应", "形成低延迟闭环", RGBColor(0xD6,0xE4,0xF7)),
              ("L3", "适应", "在线调整状态或结构", RGBColor(0xBB,0xD2,0xF2)),
              ("L4", "创造", "形成新策略或新吸引子", RGBColor(0x8F,0xB6,0xEC)),
              ("L5", "通用", "跨环境迁移与整合", MIDB),
              ("L6", "超级", "受控自演化与能力提升", NAVY)]
    lx = 4.95; y0 = CTOP + 0.15; rh = (CBOT - CTOP - 0.6) / 6; gap = 0.12
    for i, (lv, n, d, bg) in enumerate(levels):
        yy = y0 + i * (rh + gap)
        white_txt = (bg == NAVY or bg == MIDB)
        add_rect(s, lx, yy, 7.7, rh, fill=bg, rounded=True)
        add_text(s, lx + 0.2, yy, 1.0, rh, [[(lv, 20, WHITE if white_txt else NAVY, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, lx + 1.3, yy, 1.5, rh, [[(n, 19, WHITE if white_txt else NAVY, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, lx + 2.9, yy, 4.6, rh, [[(d, 15, WHITE if white_txt else TEXT, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 11)

def s12_tcc(prs):
    s = blank(prs); header(s, "用液态拓扑实现复杂性匹配 · TCC", "度量与匹配 · 12")
    add_rect(s, 0.55, CTOP + 0.1, 4.0, CBOT - CTOP - 0.2, fill=NAVY, rounded=False)
    add_text(s, 0.85, CTOP + 0.4, 3.4, 0.4, [[("TOPOLOGY-CENTRIC", 13, RGBColor(0x9F,0xB6,0xE0), False, False)]])
    add_text(s, 0.85, CTOP + 1.0, 3.4, 1.0, [[("TCC", 34, WHITE, True, False)], [("拓扑中心计算", 22, RGBColor(0xC9,0xD8,0xF0), False, False)]])
    add_rect(s, 0.85, CBOT - 1.7, 3.4, 1.4, fill=RGBColor(0x16,0x4A,0x9C), rounded=True)
    add_text(s, 1.05, CBOT - 1.55, 3.0, 1.1,
             [[("第一阶段目标：让系统复杂度", 13, RGBColor(0xEA,0xF1,0xFB), False, False),
               ("动态匹配", 14, YEL, True, False), ("环境复杂度。", 13, RGBColor(0xEA,0xF1,0xFB), False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    lx = 4.95
    add_text(s, lx, CTOP + 0.2, 7.7, 0.6, [[("系统复杂度匹配环境复杂度", 22, BLUE, True, False)]])
    add_rect(s, lx, CTOP + 0.85, 7.7, 0.03, fill=LINE, rounded=False)
    add_text(s, lx, CTOP + 1.1, 7.7, 1.7,
             [[("TCC 基于 SDI，从", 17, TEXT, False, False),
               ("复杂度视角", 17, NAVY, True, False),
               ("调节拓扑，使系统复杂度动态逼近环境复杂度。当结构被持续塑形，资源中的非线性被唤醒，产生系统级超线性增益。", 17, TEXT, False, False)]], space_after=8)
    fp = render_formula(r"C_{\mathrm{ST}}^{\mathrm{sys}} \approx C_{\mathrm{ST}}^{\mathrm{env}} \quad \Rightarrow \quad 1+1>2", "tcc", fontsize=26, color="#0E3F8C")
    add_rect(s, lx, CTOP + 3.0, 7.7, 1.2, fill=LBLUE2, line=MIDB, line_w=1.0)
    add_pic(s, fp, lx + 1.0, CTOP + 3.12, width=5.7)
    add_text(s, lx, CBOT - 0.5, 7.7, 0.4,
             [[("页脚佐证：Dally & Towles, Interconnection Networks；Hennessy & Patterson, 2019。", 11, AUX, False, False)]])
    footer(s, 12)

def s13_route(prs):
    s = blank(prs); header(s, "路径成算：Route ≈ Transform", "度量与匹配 · 13")
    lx = 0.55
    add_text(s, lx, CTOP + 0.2, 7.0, 0.6, [[("数据走过什么拓扑，就经历什么变换", 20, BLUE, True, False)]])
    add_rect(s, lx, CTOP + 0.85, 7.0, 0.03, fill=LINE, rounded=False)
    add_text(s, lx, CTOP + 1.1, 7.0, 1.8,
             [[("在可重构拓扑中，路径上的", 16, TEXT, False, False),
               ("边权、时延、节点状态、非线性器件与记忆效应", 16, NAVY, True, False),
               ("共同构成一个复合算子：数据所经历的变换，等价于它所走过的拓扑。", 16, TEXT, False, False)]], space_after=8)
    add_rect(s, lx, CTOP + 3.2, 7.0, 1.1, fill=YELBG, line=YEL, line_w=1.2)
    add_text(s, lx + 0.2, CTOP + 3.28, 6.6, 0.95,
             [[("严谨口径：Route ≈ Transform 不是无条件等号，而是拓扑可编译条件下的近似等价。", 14, RGBColor(0x8A,0x6D,0x2B), False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, lx, CBOT - 0.5, 7.0, 0.4,
             [[("页脚佐证：Maass et al., Neural Computation, 2002；Jaeger & Haas, Science, 2004。", 11, AUX, False, False)]])
    # right diagram
    add_rect(s, 7.85, CTOP + 0.1, 4.9, CBOT - CTOP - 0.2, fill=LBLUE2, line=LINE, line_w=1.0)
    add_text(s, 8.05, CTOP + 0.35, 4.5, 0.4, [[("路径即复合变换算子", 15, NAVY, True, False)]], align=PP_ALIGN.CENTER)
    fpp = [render_formula(r"\phi_{e_1}", "phi1", fontsize=22, color="#0E3F8C"),
           render_formula(r"\phi_{e_2}", "phi2", fontsize=22, color="#0E3F8C"),
           render_formula(r"\cdots", "cdots", fontsize=22, color="#0E3F8C"),
           render_formula(r"\phi_{e_k}", "phiK", fontsize=22, color="#0E3F8C")]
    bw = 1.0; gap = 0.25; x0 = 8.05; y = CTOP + 1.2
    for i, p in enumerate(fpp):
        x = x0 + i * (bw + gap)
        add_rect(s, x, y, bw, 0.9, fill=WHITE, line=MIDB, line_w=1.2, rounded=True)
        add_pic(s, p, x + bw/2 - 0.35, y + 0.18, width=0.7)
        if i < 3:
            add_text(s, x + bw, y, gap, 0.9, [[("∘", 18, BLUE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    fpT = render_formula(r"\mathcal{T}_{\rho}", "T_rho", fontsize=26, color="#FFFFFF")
    add_text(s, 8.05, y + 1.2, 4.5, 0.4, [[("=", 22, BLUE, True, False)], [("  复合算子", 15, SUB, False, False)]], align=PP_ALIGN.CENTER)
    add_rect(s, 9.0, y + 1.7, 2.6, 0.9, fill=NAVY, rounded=True)
    add_pic(s, fpT, 9.0 + 1.3 - 0.65, y + 1.88, width=1.3)
    footer(s, 13)

def s14_cmcs(prs):
    s = blank(prs); header(s, "从复杂性匹配到复杂性同步", "度量与匹配 · 14")
    add_text(s, 0, CTOP + 0.15, SW, 0.5, [[("CM 是匹配，CS 是同步，EI 是涌现", 22, BLUE, True, False)]], align=PP_ALIGN.CENTER)
    stages = [("CM", "复杂性匹配", LBLUE, NAVY), ("CS", "复杂性同步", RGBColor(0xBB,0xD2,0xF2), NAVY), ("EI", "复杂性涌现", NAVY, WHITE)]
    bw = 3.1; gap = 1.0; x0 = (SW - (3*bw + 2*gap)) / 2; y = CTOP + 0.9
    for i, (t, d, bg, fg) in enumerate(stages):
        x = x0 + i * (bw + gap)
        add_rect(s, x, y, bw, 1.6, fill=bg, rounded=True)
        add_text(s, x, y + 0.25, bw, 0.8, [[(t, 34, fg, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + 1.05, bw, 0.45, [[(d, 18, fg, True, False)]], align=PP_ALIGN.CENTER)
        if i < 2:
            add_text(s, x + bw, y, gap, 1.6, [[("→", 30, BLUE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cards = [("CM：匹配", "TCC 使系统复杂度匹配环境，形成 1+1>2 的超线性增益。", LINE, BLUE),
             ("CS：同步", "多子系统高阶复杂特征发生同步，迈向 1+1>N。", MIDB, NAVY),
             ("EI：涌现", "智能作为系统演化的涌现属性显现。", NAVY, WHITE)]
    y2 = CTOP + 3.0; ch = 2.0
    for i, (t, d, top, fg) in enumerate(cards):
        x = x0 + i * (bw + gap)
        add_rect(s, x, y2, bw, ch, fill=(NAVY if top == NAVY else LBLUE2), line=None if top == NAVY else LINE, line_w=1.0)
        add_rect(s, x, y2, bw, 0.12, fill=top, rounded=False)
        add_text(s, x, y2 + 0.3, bw, 0.6, [[(t, 20, fg, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.2, y2 + 1.0, bw - 0.4, ch - 1.1, [[(d, 15, fg if top == NAVY else TEXT, False, False)]], align=PP_ALIGN.CENTER)
    add_text(s, 0, CBOT + 0.0, SW, 0.4,
             [[("术语统一：CS = Complexity Synchronization。页脚佐证：Mahmoodi, Kerick & West, Scientific Reports, 2024。", 11, AUX, False, False)]], align=PP_ALIGN.CENTER)
    footer(s, 14)

def s15_sdde(prs):
    s = blank(prs); header(s, "先有演化语言，再谈涌现机制 · SDDE", "演化与调控 · 15")
    add_rect(s, 0.55, CTOP + 0.1, 4.0, CBOT - CTOP - 0.2, fill=NAVY, rounded=False)
    add_text(s, 0.85, CTOP + 0.4, 3.4, 0.4, [[("STOCHASTIC DELAY DE", 13, RGBColor(0x9F,0xB6,0xE0), False, False)]])
    add_text(s, 0.85, CTOP + 1.0, 3.4, 1.0, [[("SDDE", 40, WHITE, True, False)], [("连续演化语言", 22, RGBColor(0xC9,0xD8,0xF0), False, False)]])
    add_rect(s, 0.85, CBOT - 1.7, 3.4, 1.4, fill=RGBColor(0x16,0x4A,0x9C), rounded=True)
    add_text(s, 1.05, CBOT - 1.55, 3.0, 1.1,
             [[("先有演化语言，再有涌现调控机制——故置于 iNEST 之前更自洽。", 14, RGBColor(0xEA,0xF1,0xFB), False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    lx = 4.95
    add_text(s, lx, CTOP + 0.2, 7.7, 0.6, [[("带噪声、带时延、带记忆的连续演化", 20, BLUE, True, False)]])
    fp = render_formula(
        r"dx_i(t) = \left[ f_i + \sum_j A_{ij}(t) w_{ij}(t) g_{ij} + u_i \right] dt + \sigma_i dW_i(t)",
        "sdde", fontsize=20, color="#0E3F8C")
    add_rect(s, lx, CTOP + 1.0, 7.7, 1.5, fill=LBLUE2, line=MIDB, line_w=1.0)
    add_pic(s, fp, lx + 0.4, CTOP + 1.2, width=6.9)
    add_text(s, lx, CTOP + 2.85, 7.7, 2.0,
             [[("SDDE 刻画", 16, TEXT, False, False), ("材料内部状态、非线性器件、拓扑耦合与环境反馈", 16, NAVY, True, False),
               ("的共同演化；其中 A_{ij}(t) 正是 SDI 可直接调控的耦合结构。", 16, TEXT, False, False)]], space_after=8)
    add_text(s, lx, CBOT - 0.5, 7.7, 0.4,
             [[("页脚佐证：Mohammed, 1984；Buckwar, 2000；Mao, 2007。", 11, AUX, False, False)]])
    footer(s, 15)

def s16_sdi(prs):
    s = blank(prs); header(s, "连续演化，需要离散拓扑执行器 · SDI", "演化与调控 · 16")
    add_text(s, 0, CTOP + 0.1, SW, 0.4,
             [[("SDI TOPOLOGY CONTROL OPERATOR（工程控制律）", 13, AUX, False, False)]], align=PP_ALIGN.CENTER)
    fp = render_formula(
        r"A(t^+) = \Pi_{\mathrm{SDI}}\left[ A(t), \lambda_{\max}^{FT}, \hat{m}, TE, AIS, E_{\mathrm{diss}}, \Gamma_{\mathrm{st}}^{u} \right]",
        "sdi", fontsize=22, color="#0E3F8C")
    add_rect(s, 0.7, CTOP + 0.6, 11.9, 1.55, fill=LBLUE2, line=BLUE, line_w=2.0)
    add_pic(s, fp, 1.4, CTOP + 0.75, width=10.5)
    # 7 param chips
    chips = [("A(t)", "当前拓扑"), ("λ_max^FT", "有限时间李雅普诺夫指数"), ("m̂", "分支比"),
             ("TE", "传递熵"), ("AIS", "动态信息存储"), ("E_diss", "耗散能量"), ("Γ_st^u", "可用时空协同因子")]
    cw = 1.6; gap = 0.12; x0 = 0.7; y = CTOP + 2.45
    for i, (k, d) in enumerate(chips):
        kk = k.replace("^", "").replace("λ_maxFT", "lambda_max").replace("m̂", "m_hat").replace("Γ_st^u", "Gamma_st")
        fps = render_formula(k.replace("λ_max^FT", "\\lambda_{\\max}^{FT}").replace("m̂", "\\hat{m}").replace("Γ_st^u", "\\Gamma_{\\mathrm{st}}^{u}").replace("E_diss", "E_{\\mathrm{diss}}"),
                              "chip_" + str(i) + "_" + kk, fontsize=18, color="#1E4FA8")
        x = x0 + i * (cw + gap)
        add_rect(s, x, y, cw, 1.15, fill=WHITE, line=LINE, line_w=1.0)
        add_pic(s, fps, x + cw/2 - 0.45, y + 0.08, width=0.9)
        add_text(s, x, y + 0.62, cw, 0.45, [[(d, 10.5, AUX, False, False)]], align=PP_ALIGN.CENTER)
    add_rect(s, 0.7, CTOP + 3.95, 11.9, 0.95, fill=BLUEBG, line=BLUE, line_w=1.2)
    add_text(s, 0.9, CTOP + 4.0, 11.5, 0.85,
             [[("SDI 拓扑调控算子：从“受控 SDDE 混合动力系统”中抽象出的", 15, NAVY, True, False),
               ("工程控制律", 16, NAVY, True, False), ("，而非自然定律公式。", 15, NAVY, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 0.7, CBOT - 1.0, 11.9, 0.85, fill=YELBG, line=YEL, line_w=1.2)
    add_text(s, 0.9, CBOT - 0.95, 11.5, 0.75,
             [[("口播边界：这是工程控制抽象，不宣称为已证明的普适定律。", 14, RGBColor(0x8A,0x6D,0x2B), False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0, CBOT + 0.02, SW, 0.32,
             [[("页脚佐证：Liu, Slotine & Barabási, Nature, 2011；Schreiber, PRL, 2000；Lizier, 2012。", 10.5, AUX, False, False)]], align=PP_ALIGN.CENTER)
    footer(s, 16)

def s17_inest(prs):
    s = blank(prs); header(s, "局部规则推动系统接近临界 · iNEST", "演化与调控 · 17")
    add_rect(s, 0.55, CTOP + 0.1, 4.0, CBOT - CTOP - 0.2, fill=NAVY, rounded=False)
    add_text(s, 0.85, CTOP + 0.4, 3.4, 0.4, [[("CRITICAL INTELLIGENCE", 12, RGBColor(0x9F,0xB6,0xE0), False, False)]])
    add_text(s, 0.85, CTOP + 1.0, 3.4, 1.0, [[("iNEST", 40, WHITE, True, False)], [("临界智能区", 22, RGBColor(0xC9,0xD8,0xF0), False, False)]])
    add_rect(s, 0.85, CBOT - 1.7, 3.4, 1.4, fill=RGBColor(0x16,0x4A,0x9C), rounded=True)
    add_text(s, 1.05, CBOT - 1.55, 3.0, 1.1,
             [[("塑边 · 选向 · 稳态 · 临界——逼近混沌边缘。", 14, RGBColor(0xEA,0xF1,0xFB), False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    lx = 4.95
    add_text(s, lx, CTOP + 0.25, 7.7, 0.9,
             [[("iNEST 通过", 16, TEXT, False, False), ("局部规则与 SDI 调控", 16, NAVY, True, False),
               ("，使复杂度在介观物理网络中持续演化，接近混沌边缘并形成稳定吸引子。", 16, TEXT, False, False)]])
    data = [("机制", "作用"),
            ("STDP", "按时序塑造连接"),
            ("局部预测误差", "引导状态向低误差吸引子收敛"),
            ("稳态可塑性", "防止全静默或全饱和"),
            ("SOC / EOC", "维持临界窗口")]
    add_table(s, lx, CTOP + 1.5, 7.7, 3.6, data, col_widths=[2.2, 5.5], font_size=15, hdr_size=16)
    add_text(s, lx, CBOT - 0.5, 7.7, 0.4,
             [[("页脚佐证：Bi & Poo, J Neurosci, 1998；Bak et al., PRL, 1987；Kinouchi & Copelli, Nature Physics, 2006。", 10.5, AUX, False, False)]])
    footer(s, 17)

def s18_imeso(prs):
    s = blank(prs); header(s, "iMESO 介观物理智能平台", "平台与路线 · 18")
    add_rect(s, 0.55, CTOP + 0.1, 6.3, CBOT - CTOP - 0.2, fill=LBLUE2, line=LINE, line_w=1.0)
    add_text(s, 0.85, CTOP + 0.35, 5.7, 0.4,
             [[("材料·器件·拓扑·动力学·环境 闭环一体化", 14, NAVY, True, False)]], align=PP_ALIGN.CENTER)
    layers = [("材料", "状态记忆", LBLUE), ("器件", "非线性载荷", RGBColor(0xD6,0xE4,0xF7)),
              ("拓扑", "协同关系", RGBColor(0xBB,0xD2,0xF2)), ("动力学", "临界演化", RGBColor(0x8F,0xB6,0xEC)),
              ("环境", "选择智能行为", NAVY)]
    y = CTOP + 1.0; lh = 0.82
    for i, (t, d, bg) in enumerate(layers):
        white = (bg == NAVY or bg == RGBColor(0x8F,0xB6,0xEC))
        add_rect(s, 0.85, y, 5.7, lh - 0.1, fill=bg, rounded=True)
        add_text(s, 1.05, y, 1.6, lh - 0.1, [[(t, 17, WHITE if white else NAVY, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.7, y, 3.7, lh - 0.1, [[(d, 15, WHITE if white else TEXT, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
        if i < 4:
            add_text(s, 0.85, y + lh - 0.1, 5.7, 0.25, [[("↑ 反馈闭环 ↓", 11, MIDB, False, False)]], align=PP_ALIGN.CENTER)
        y += lh + 0.15
    lx = 7.25
    add_text(s, lx, CTOP + 0.3, 5.5, 0.6, [[("让智能涌现走向工程验证", 20, BLUE, True, False)]])
    add_rect(s, lx, CTOP + 0.95, 5.5, 0.03, fill=LINE, rounded=False)
    add_text(s, lx, CTOP + 1.2, 5.5, 1.8,
             [[("iMESO 作为", 16, TEXT, False, False), ("晶圆 / 晶矩 / 面板级", 16, NAVY, True, False),
               ("介观物理载体，把材料、器件、拓扑、动力学、环境闭环集成于同一平台，使前述理论可被工程验证。", 16, TEXT, False, False)]], space_after=8)
    add_rect(s, lx, CTOP + 3.2, 5.5, 1.4, fill=BLUEBG, line=BLUE, line_w=1.2)
    add_text(s, lx + 0.2, CTOP + 3.3, 5.1, 1.2,
             [[("材料提供状态记忆，拓扑组织协同关系，临界放大有效响应，环境选择智能行为。", 16, NAVY, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, lx, CBOT - 0.5, 5.5, 0.4,
             [[("页脚佐证：Strukov et al., 2008；Marković et al., 2020；Sebastian et al., Nat Nano, 2020；Schuman et al., Nat Comput Sci, 2022。", 10.5, AUX, False, False)]])
    footer(s, 18)

def s19_roadmap(prs):
    s = blank(prs); header(s, "从百万连接到百万亿连接 · 智涌脑路线图", "平台与路线 · 19")
    add_text(s, 0.55, CTOP + 0.15, 12.0, 0.5, [[("每两年网络规模提升 100 倍", 20, BLUE, True, False)]])
    data = [("时间", "网络连接规模", "智涌脑阶段", "智能等级目标"),
            ("2027", "10⁶ 百万级", "原理验证智涌脑", "感知 L1"),
            ("2029", "10⁸ 亿级", "端侧智涌脑", "反应 L2"),
            ("2031", "10¹⁰ 百亿级", "边侧智涌脑", "适应 L3"),
            ("2033", "10¹² 万亿级", "创造型智涌脑", "创造 L4"),
            ("2035", "10¹⁴ 百万亿级", "类人脑规模智涌脑", "通用智能 L5")]
    add_table(s, 0.55, CTOP + 0.8, 12.2, 3.5, data, col_widths=[1.4, 3.0, 4.0, 3.8], font_size=16, hdr_size=17)
    add_rect(s, 0.55, CTOP + 4.5, 12.2, 1.1, fill=YELBG, line=YEL, line_w=1.2)
    add_text(s, 0.75, CTOP + 4.58, 11.8, 0.95,
             [[("边界口径：连接规模不是智能本身，而是承载复杂度演化的物理容量；智能等级取决于 CST 相对环境复杂度的裕度（见 P10–P11）。", 15, RGBColor(0x8A,0x6D,0x2B), False, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.55, CBOT - 0.05, 12.2, 0.4,
             [[("页脚佐证：Azevedo et al., J Comp Neurol, 2009；Sporns, Networks of the Brain, 2011。", 11, AUX, False, False)]])
    footer(s, 19)

def s20_ending(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SW, SH, fill=NAVY, rounded=False)
    add_rect(s, 0, 0, 0.35, SH, fill=MIDB, rounded=False)
    add_text(s, 0, 0.6, SW, 0.4,
             [[("FROM SOFTWARE-DEFINED TO COMPLEXITY-DEFINED", 13, RGBColor(0x9F,0xB6,0xE0), False, False)]], align=PP_ALIGN.CENTER)
    add_text(s, 0, 1.1, SW, 1.0, [[("从软件定义系统，到复杂度定义智能", 38, WHITE, True, False)]], align=PP_ALIGN.CENTER)
    pills = ["SDSoW 给资源", "CST 给刻度", "TCC 做匹配", "SDDE 写演化", "SDI 做调控", "iNEST 推临界", "iMESO 造智涌脑"]
    bw = 2.5; gap = 0.35; x0 = (SW - (7*bw + 6*gap)) / 2; y = 2.7
    for i, p in enumerate(pills):
        x = x0 + i * (bw + gap)
        add_rect(s, x, y, bw, 0.85, fill=RGBColor(0x16,0x4A,0x9C), line=MIDB, line_w=1.2, rounded=True)
        add_text(s, x, y, bw, 0.85, [[(p, 16, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 2.2, 4.6, 8.9, 1.9, fill=RGBColor(0x16,0x4A,0x9C), line=YEL, line_w=1.5, rounded=False)
    add_text(s, 2.5, 4.75, 8.3, 1.4,
             [[("先让系统复杂度匹配环境，再让复杂度在临界演化中同步放大，最终让介观物理网络自己长出智能。", 20, YEL, True, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0, 6.6, SW, 0.4, [[("—— 邬江兴 院士", 15, RGBColor(0xDC,0xE7,0xF8), False, False)]], align=PP_ALIGN.CENTER)

# ---------------- build ----------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    builders = [s01_cover, s02_frontier, s03_walls, s04_moore, s05_sdsow, s06_nonlinear,
                s07_threshold, s08_cst, s09_exponent, s10_ri, s11_levels, s12_tcc, s13_route,
                s14_cmcs, s15_sdde, s16_sdi, s17_inest, s18_imeso, s19_roadmap, s20_ending]
    for b in builders:
        b(prs)
    out = os.path.join(OUTDIR, "介观尺度计算新范式.pptx")
    prs.save(out)
    print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    main()
